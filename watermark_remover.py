from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt
import os

class PPTXWatermarkRemover:
    def __init__(self, target_domain="gamma.app", corner_threshold=0.7):
        """
        target_domain: string contained in hyperlinks to remove (case-insensitive)
        corner_threshold: fraction of slide width/height defining 'bottom-right corner'
        """
        self.target_domain = target_domain.lower()
        self.corner_threshold = corner_threshold

        self.links_removed = 0
        self.shapes_removed = 0
        self.corner_images_removed = 0

    # ---------- Public API ----------
    def clean_pptx_from_target_domain(self, pptx_path, output_path):
        prs = Presentation(pptx_path)

        print(f"Processing file: {pptx_path}")
        print(f"Target domain: {self.target_domain}")
        print(f"Slides: {len(prs.slides)}")

        # Masters and layouts first (watermarks often live there)
        try:
            for master in prs.part.slide_masters:
                self._process_shape_tree(master, context="MASTER")
        except Exception as e:
            print(f"! Could not process masters: {e}")

        for layout in prs.slide_layouts:
            self._process_shape_tree(layout, context="LAYOUT")

        # Then actual slides
        for i, slide in enumerate(prs.slides, start=1):
            print(f"\nSlide {i}:")
            self._process_slide(slide)

        prs.save(output_path)

        print("\n" + "="*60)
        print("RESULT")
        print(f"Links removed: {self.links_removed}")
        print(f"Shapes removed (incl. pictures): {self.shapes_removed}")
        print(f"Corner pictures removed: {self.corner_images_removed}")
        print(f"Cleaned file: {output_path}")

        return self.shapes_removed, self.links_removed, self.corner_images_removed

    # ---------- Internals ----------
    def _process_slide(self, slide):
        sw, sh = self._slide_size(slide)
        right_edge = sw * self.corner_threshold
        bottom_edge = sh * self.corner_threshold

        # 1) If any linked image in corner -> remove all corner pictures
        corner_picture_ids_to_remove = set()
        corner_has_target_link = False

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and self._is_in_corner(shape, right_edge, bottom_edge):
                if self._shape_has_target_link(shape):
                    corner_has_target_link = True
                corner_picture_ids_to_remove.add(id(shape))

        if corner_has_target_link:
            removed = self._remove_shapes_by_id(slide, corner_picture_ids_to_remove)
            self.corner_images_removed += removed
            if removed:
                print(f"  ✓ Removed {removed} corner picture(s) due to target-domain link")

        # 2) Remove any shape/picture with hyperlink to target domain
        # 3) Strip text run hyperlinks to target domain
        # (Do second pass so we don’t mutate while we iterate corner set above)
        # collect to remove
        to_remove_ids = []

        for shape in list(slide.shapes):
            # already removed in step 1
            if id(shape) not in [id(s) for s in slide.shapes]:
                continue

            if self._shape_has_target_link(shape):
                to_remove_ids.append(id(shape))
                continue

            # text runs: remove hyperlink or the run if it is only a link
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                changed = self._strip_text_run_links(shape)
                if changed:
                    print(f"  ✓ Stripped {changed} hyperlink(s) from text")
                    self.links_removed += changed

        removed = self._remove_shapes_by_id(slide, to_remove_ids)
        self.shapes_removed += removed
        if removed:
            print(f"  ✓ Removed {removed} shape(s) with target-domain link")

    def _process_shape_tree(self, obj_with_shapes, context=""):
        # masters/layouts
        try:
            for shape in list(obj_with_shapes.shapes):
                # remove shapes with direct hyperlink
                if self._shape_has_target_link(shape):
                    obj_with_shapes.shapes._spTree.remove(shape._element)
                    self.shapes_removed += 1
                    print(f"  ✓ [{context}] Removed shape with target-domain link")
                    continue

                # strip text links inside masters/layouts
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    changed = self._strip_text_run_links(shape)
                    if changed:
                        print(f"  ✓ [{context}] Stripped {changed} hyperlink(s) from text")
                        self.links_removed += changed
        except Exception as e:
            print(f"  ! [{context}] Error processing shapes: {e}")

    def _remove_shapes_by_id(self, slide, ids_set):
        removed = 0
        for shape in list(slide.shapes):
            if id(shape) in ids_set:
                slide.shapes._spTree.remove(shape._element)
                removed += 1
        return removed

    def _shape_has_target_link(self, shape):
        try:
            # picture or auto-shape click hyperlink
            if hasattr(shape, "click_action") and shape.click_action is not None:
                h = shape.click_action.hyperlink
                if h is not None and h.address:
                    if self.target_domain in h.address.lower():
                        self.links_removed += 1  # count link removed with shape
                        return True

            # text-frame link on the whole shape (rare, but possible)
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.hyperlink and r.hyperlink.address:
                            if self.target_domain in r.hyperlink.address.lower():
                                # we’ll remove link at run-level instead of whole shape
                                # but for “shape_has_target_link” we leave False
                                pass
        except Exception:
            pass
        return False

    def _strip_text_run_links(self, shape):
        """Remove hyperlinks from text runs that point to target domain.
           Returns the number of links stripped.
        """
        removed = 0
        tf = shape.text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                try:
                    if r.hyperlink and r.hyperlink.address and self.target_domain in r.hyperlink.address.lower():
                        # Clear the hyperlink (python-pptx: setting address to None detaches it)
                        r.hyperlink.address = None
                        removed += 1
                        # Optionally, remove the run if it was pure link text with no value after strip
                        if r.text.strip() == "":
                            r._r.getparent().remove(r._r)
                except Exception:
                    continue
        return removed

    def _slide_size(self, slide):
        # slide width/height from presentation object
        # pptx stores sizes at presentation level; fetch via slide.part
        prs = slide.part.package.presentation_part.presentation
        return float(prs.slide_width), float(prs.slide_height)

    def _is_in_corner(self, shape, right_edge, bottom_edge):
        try:
            # left/top/width/height are EMUs
            left = float(shape.left)
            top = float(shape.top)
            return (left >= right_edge) and (top >= bottom_edge)
        except Exception:
            return False


if __name__ == "__main__":
    # Example usage
    inp = "input.pptx"
    out = "cleaned_without_gamma_links.pptx"
    remover = PPTXWatermarkRemover(target_domain="gamma.app", corner_threshold=0.7)
    shapes_removed, links_removed, corner_removed = remover.clean_pptx_from_target_domain(inp, out)